"""
master_analyzer.py
------------------
Phân tích sâu file PPTX master slide, trích xuất toàn bộ thông tin cần thiết
để inject vào prompt cho AI gen SVG phù hợp với cấu trúc master.

Thông tin trích xuất:
  - Tất cả layouts: tên, data-layout value tương ứng, danh sách placeholders
  - Mỗi placeholder: idx, type, label, tọa độ px (1280×720), font size, màu
  - Theme colors: dk1, lt1, dk2, lt2, accent1–6
  - Fonts: heading font, body font
  - Background: solid/gradient/inherit cho mỗi layout
"""

import io
from lxml import etree
from pptx import Presentation

# ─── Target canvas (normalize tất cả EMU về đây) ──────────────────────────────────
_DEFAULT_TARGET_W_PX = 1280
_DEFAULT_TARGET_H_PX = 720
# Scale được tính động trong analyze_master() từ prs.slide_width/height

# ─── Map tên layout (multi-language) → data-layout value dùng trong SVG ────────
_LAYOUT_NAME_MAP = {
    # Tiếng Anh
    "title slide":                      "title-slide",
    "title and content":                "content",
    "two content":                      "two-column",
    "comparison":                       "comparison",
    "section header":                   "section-header",
    "blank":                            "blank",
    "content with caption":             "content-caption",
    "picture with caption":             "picture-caption",
    "title only":                       "title-only",

    # Tiếng Nhật (VTI master ver1 + ver2)
    "タイトル スライド":                 "title-slide",
    "タイトルとコンテンツ":              "content",
    "セクション見出し":                  "section-header",
    "2 つのコンテンツ":                  "two-column",
    "比較":                              "comparison",
    "タイトルのみ":                      "title-only",
    "タイトル付きのコンテンツ":          "content-caption",
    "タイトル付きの図":                  "picture-caption",
    "タイトルと縦書きテキスト":          "content",
    "縦書きタイトルと\n縦書きテキスト":  "content",
    "白紙":                              "blank",

    # Korean
    "제목 슬라이드":                     "title-slide",
    "제목 및 내용":                      "content",
    "두 내용":                           "two-column",
    "구역 머리글":                    "section-header",

    # Chinese
    "标题幻灯片":                        "title-slide",
    "标题和内容":                        "content",
    "两栏内容":                          "two-column",
    "节标题":                            "section-header",
}

# ─── Japanese / alias font fallback chains ──────────────────────────────────
_FONT_FALLBACK_MAP = {
    "游ゴシック light":   "Yu Gothic Light, Meiryo, Arial, sans-serif",
    "yu gothic light":    "Yu Gothic Light, Meiryo, Arial, sans-serif",
    "游ゴシック":          "Yu Gothic, Meiryo, Arial, sans-serif",
    "yu gothic":          "Yu Gothic, Meiryo, Arial, sans-serif",
    "游明朝":              "Yu Mincho, Times New Roman, serif",
    "メイリオ":            "Meiryo, Arial, sans-serif",
    "meiryo":             "Meiryo, Arial, sans-serif",
    "+mj-ea":             "Yu Gothic Light, Meiryo, sans-serif",
    "+mn-ea":             "Yu Gothic, Meiryo, sans-serif",
    "+mj-lt":             "Calibri Light, Arial, sans-serif",
    "+mn-lt":             "Calibri, Arial, sans-serif",
}


def _resolve_font(font_name: str) -> str:
    """Thêm fallback chain cho Japanese/alias fonts."""
    return _FONT_FALLBACK_MAP.get(font_name.lower().strip(), font_name)

# ─── Map placeholder idx → type string ──────────────────────────────────────
_PH_IDX_TYPE_MAP = {
    0:  "center_title",
    1:  "body",
    2:  "subtitle",
    3:  "title",
    10: "date",
    11: "footer",
    12: "slide_number",
    13: "object",
    14: "picture",
}


def _emu_to_px(emu_val: int, scale: float) -> int:
    """Chuyển đổi EMU → px dùng scale động từ file thực tế."""
    return round(int(emu_val or 0) * scale)


def _extract_theme_data(prs: Presentation) -> dict:
    """
    Đọc theme colors và font names từ XML nội bộ của slide master.
    Trả về dict với keys: fonts {heading, body, heading_raw, body_raw},
    colors {dk1, lt1, ...accent1-6}.
    """
    fonts  = {"heading": "Calibri Light", "body": "Calibri",
              "heading_raw": "", "body_raw": ""}
    colors = {}

    try:
        master_part = prs.slide_master.part
        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"

        for rel in master_part.rels.values():
            if "theme" not in rel.reltype:
                continue
            try:
                theme_xml = etree.fromstring(rel.target_part.blob)
            except Exception:
                continue

            ns = {"a": ns_a}

            # ── Fonts ──────────────────────────────────────────────────────
            try:
                major = theme_xml.find(".//a:fontScheme/a:majorFont/a:latin", ns)
                minor = theme_xml.find(".//a:fontScheme/a:minorFont/a:latin", ns)
                if major is not None and major.get("typeface"):
                    raw = major.get("typeface")
                    fonts["heading_raw"] = raw
                    fonts["heading"]     = _resolve_font(raw)
                if minor is not None and minor.get("typeface"):
                    raw = minor.get("typeface")
                    fonts["body_raw"] = raw
                    fonts["body"]     = _resolve_font(raw)
            except Exception:
                pass

            # ── Colors ─────────────────────────────────────────────────────
            try:
                wanted = {"dk1", "lt1", "dk2", "lt2",
                          "accent1", "accent2", "accent3",
                          "accent4", "accent5", "accent6"}
                for node in theme_xml.findall(".//a:clrScheme/*", ns):
                    try:
                        tag = etree.QName(node.tag).localname
                        if tag not in wanted:
                            continue
                        child = list(node)
                        if not child:
                            continue
                        hex_val = (child[0].get("val") or child[0].get("lastClr", "")).upper()
                        if hex_val:
                            colors[tag] = f"#{hex_val.lstrip('#')}"
                    except Exception:
                        continue
            except Exception:
                pass

            break  # chỉ cần theme đầu tiên

    except Exception:
        pass  # fallback về giá trị mặc định

    return {"fonts": fonts, "colors": colors}


def _extract_background(layout) -> dict:
    """
    Trích xuất thông tin background của một layout.
    Trả về {"type": "solid"|"gradient"|"inherit", "colors": [hex...]}.
    Tất cả exception được catch silently.
    """
    try:
        fill = layout.background.fill
        if fill.type is not None:
            fill_type_str = str(fill.type)
            if "SOLID" in fill_type_str or fill_type_str == "1":
                try:
                    hex_color = f"#{fill.fore_color.rgb}"
                    return {"type": "solid", "colors": [hex_color]}
                except Exception:
                    pass
    except Exception:
        pass

    try:
        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        grad = layout.element.find(f".//{{{ns_a}}}gradFill")
        if grad is not None:
            stops = grad.findall(f".//{{{ns_a}}}gs")
            hex_colors = []
            for stop in stops:
                try:
                    srgb = stop.find(f"{{{ns_a}}}srgbClr")
                    if srgb is not None:
                        val = srgb.get("val", "")
                        if val:
                            hex_colors.append(f"#{val.upper()}")
                except Exception:
                    continue
            if hex_colors:
                return {"type": "gradient", "colors": hex_colors}
    except Exception:
        pass

    return {"type": "inherit", "colors": []}


def _extract_bullet_levels(ph) -> list:
    """
    Đọc font size theo từng cấp indent trong body placeholder.
    Trả về list[{"level": int, "font_size_pt": int, "indent_px": int}].
    """
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    levels = []

    try:
        tx_body = ph._element.find(f"{{{ns_a}}}txBody")
        if tx_body is None:
            raise ValueError("no txBody")

        for i in range(1, 4):
            tag   = f"lvl{i}pPr"
            sz_pt = 18
            try:
                node = tx_body.find(f".//{{{ns_a}}}{tag}")
                if node is not None:
                    def_rpr = node.find(f"{{{ns_a}}}defRPr")
                    if def_rpr is not None:
                        sz_raw = def_rpr.get("sz")
                        if sz_raw and sz_raw.isdigit():
                            sz_pt = int(sz_raw) // 100
            except Exception:
                pass
            levels.append({
                "level":        i,
                "font_size_pt": sz_pt,
                "indent_px":    (i - 1) * 32,
            })
    except Exception:
        levels = [
            {"level": 1, "font_size_pt": 18, "indent_px": 0},
            {"level": 2, "font_size_pt": 16, "indent_px": 32},
            {"level": 3, "font_size_pt": 14, "indent_px": 64},
        ]

    return levels


def _extract_ph_font_info(ph) -> dict:
    """
    Đọc font từ lstStyle/lvl1pPr/defRPr trong XML của placeholder.
    Đây là cách đúng cho master/layout placeholders (không dùng text_frame.runs
    vì master placeholders thường không có runs).

    VTI ver2 verified values:
      Layout[1] title:    Meiryo 36pt BOLD #0050AD
      Layout[0] title:    Meiryo 60pt BOLD #0050AD align=ctr
      Layout[0] subtitle: Meiryo 24pt align=ctr
      slide_number:       Meiryo 16pt BOLD color=scheme:bg1
    """
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    result = {
        "font_size_pt": 18,
        "font_bold":    False,
        "font_color":   "",
        "font_name":    "",
        "align":        "",
    }
    try:
        txBody = ph._element.find(f"{{{ns_a}}}txBody")
        if txBody is None:
            return result
        lstStyle = txBody.find(f"{{{ns_a}}}lstStyle")
        if lstStyle is None:
            return result
        lvl1 = lstStyle.find(f"{{{ns_a}}}lvl1pPr")
        if lvl1 is None:
            return result

        result["align"] = lvl1.get("algn", "")

        defRPr = lvl1.find(f"{{{ns_a}}}defRPr")
        if defRPr is None:
            return result

        sz = defRPr.get("sz")
        if sz and sz.isdigit():
            result["font_size_pt"] = int(sz) // 100

        result["font_bold"] = defRPr.get("b") == "1"

        latin = defRPr.find(f"{{{ns_a}}}latin")
        if latin is not None:
            result["font_name"] = _resolve_font(latin.get("typeface", ""))

        srgb = defRPr.find(f".//{{{ns_a}}}srgbClr")
        if srgb is not None:
            result["font_color"] = f"#{srgb.get('val', '').upper()}"
        else:
            schClr = defRPr.find(f".//{{{ns_a}}}schemeClr")
            if schClr is not None:
                result["font_color"] = f"scheme:{schClr.get('val', '')}"

    except Exception:
        pass
    return result


def _calc_content_zone(
    placeholders: list,
    scale_x: float,
    scale_y: float,
    w_px: int,
    h_px: int,
) -> dict:
    """
    Tính vùng content tự do dựa trên các placeholder đã có.
    Logic: content zone = phần slide không bị chiếm bởi title/footer/slide_number.

    Trả về {"x": int, "y": int, "w": int, "h": int} (px, canvas base).
    """
    MARGIN        = 8
    header_bottom = MARGIN
    footer_top    = h_px - MARGIN

    for ph in placeholders:
        ph_type = ph.get("type", "")
        t = ph.get("top_px", 0)
        h = ph.get("height_px", 0)
        b = t + h

        if ph_type in ("center_title", "title", "date") and t < h_px // 2:
            if b + MARGIN > header_bottom:
                header_bottom = b + MARGIN

        if ph_type in ("slide_number", "footer") and t > h_px // 2:
            if t - MARGIN < footer_top:
                footer_top = t - MARGIN

    left  = round(8 / 1280 * w_px)
    right = w_px - left

    return {
        "x": left,
        "y": header_bottom,
        "w": right - left,
        "h": max(0, footer_top - header_bottom),
    }


def analyze_master(pptx_bytes: bytes) -> dict:
    """
    Hàm chính: phân tích PPTX và trả về master schema đầy đủ.
    Tất cả exception nội bộ được catch silently → luôn trả về dict hợp lệ.

    Returns dict:
    {
      "meta": {
        "slide_width_emu": int, "slide_height_emu": int,
        "slide_width_px": 1280, "slide_height_px": 720,
        "emu_scale_x": float, "emu_scale_y": float
      },
      "theme": {
        "fonts": {"heading": str, "body": str, "heading_raw": str, "body_raw": str},
        "colors": {"dk1": "#hex", ...}
      },
      "layouts": [
        {
          "index": int, "name": str, "data_layout_value": str,
          "background": {...},
          "content_zone": {"x": int, "y": int, "w": int, "h": int},
          "placeholders": [
            {
              "idx": int, "type": str, "label": str,
              "left_px": int, "top_px": int, "width_px": int, "height_px": int,
              "left_emu": int, "top_emu": int, "width_emu": int, "height_emu": int,
              "font_size_pt": int, "font_bold": bool, "font_color": str,
              "font_name": str, "align": str,
              "bullet_levels": [...]   # chỉ có khi type == "body"
            }
          ]
        }
      ]
    }
    """
    try:
        prs = Presentation(io.BytesIO(pptx_bytes))
    except Exception as exc:
        raise ValueError(f"Không thể đọc file PPTX: {exc}") from exc

    # ── Scale động theo kích thước thực tế ──────────────────────────────────
    _w_emu   = int(prs.slide_width)
    _h_emu   = int(prs.slide_height)
    _scale_x = _DEFAULT_TARGET_W_PX / _w_emu
    _scale_y = _DEFAULT_TARGET_H_PX / _h_emu

    schema = {
        "meta": {
            "slide_width_emu":  _w_emu,
            "slide_height_emu": _h_emu,
            "slide_width_px":   _DEFAULT_TARGET_W_PX,
            "slide_height_px":  _DEFAULT_TARGET_H_PX,
            "emu_scale_x":      _scale_x,
            "emu_scale_y":      _scale_y,
        },
        "theme":   _extract_theme_data(prs),
        "layouts": [],
    }

    for idx, layout in enumerate(prs.slide_layouts):
        try:
            name        = (layout.name or f"Layout {idx}").strip()
            data_layout = _LAYOUT_NAME_MAP.get(name.lower(), "content")

            layout_info = {
                "index":             idx,
                "name":              name,
                "data_layout_value": data_layout,
                "background":        _extract_background(layout),
                "placeholders":      [],
            }

            for ph in layout.placeholders:
                try:
                    ph_idx  = ph.placeholder_format.idx
                    ph_type = _PH_IDX_TYPE_MAP.get(ph_idx, "body")

                    # ── Tọa độ ───────────────────────────────────────────────────
                    try:
                        left   = int(ph.left   or 0)
                        top    = int(ph.top    or 0)
                        width  = int(ph.width  or 0)
                        height = int(ph.height or 0)
                    except Exception:
                        left = top = width = height = 0

                    # ── Font từ lstStyle XML ──────────────────────────────────────
                    fi           = _extract_ph_font_info(ph)
                    font_size_pt = fi["font_size_pt"]
                    font_bold    = fi["font_bold"]
                    font_color   = fi["font_color"]
                    font_name    = fi["font_name"]
                    align        = fi["align"]

                    ph_data = {
                        "idx":          ph_idx,
                        "type":         ph_type,
                        "label":        ph.name or f"ph_{ph_idx}",
                        "left_emu":     left,  "top_emu":    top,
                        "width_emu":    width, "height_emu": height,
                        "left_px":      _emu_to_px(left,   _scale_x),
                        "top_px":       _emu_to_px(top,    _scale_y),
                        "width_px":     _emu_to_px(width,  _scale_x),
                        "height_px":    _emu_to_px(height, _scale_y),
                        "font_size_pt": font_size_pt,
                        "font_bold":    font_bold,
                        "font_color":   font_color,
                        "font_name":    font_name,
                        "align":        align,
                    }

                    if ph_type == "body":
                        ph_data["bullet_levels"] = _extract_bullet_levels(ph)

                    layout_info["placeholders"].append(ph_data)

                except Exception:
                    continue  # bỏ qua placeholder lỗi, tiếp tục

            # ── Content zone (vùng AI tự do design) ─────────────────────────────
            layout_info["content_zone"] = _calc_content_zone(
                layout_info["placeholders"], _scale_x, _scale_y,
                _DEFAULT_TARGET_W_PX, _DEFAULT_TARGET_H_PX
            )

            schema["layouts"].append(layout_info)

        except Exception:
            continue  # bỏ qua layout lỗi, tiếp tục

    return schema
