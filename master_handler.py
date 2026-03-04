"""
master_handler.py
-----------------
Xử lý master slide từ file PPTX upload.
Đọc thông tin layout, theme, placeholder để map SVG semantic content vào đúng vị trí.

Các hàm chính:
  - parse_master_info(pptx_bytes)         → dict chứa layout, theme, font
  - find_best_layout(master_info, layout) → index layout phù hợp nhất
  - extract_svg_semantic_content(svg)     → dict nội dung từ data-role attributes
  - build_pptx_with_master(slides, ...)   → io.BytesIO file PPTX đã map vào master
"""

import io
import re
from lxml import etree
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN


# ─── Namespace cho DrawingML theme XML ───────────────────────────────────────
_DML_NS  = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
_SVG_NS  = "http://www.w3.org/2000/svg"

# ─── SVG canvas + content zone (pixel, tương ứng viewBox="0 0 1280 720") ────
_SVG_W   = 1280   # SVG canvas width  (px)
_SVG_H   = 720    # SVG canvas height (px)
_CZ_X    = 88     # content zone left  (px) — bắt đầu từ margin trái
_CZ_Y    = 131    # content zone top   (px) — kết thúc header zone
_CZ_W    = 1104   # content zone width (px)
_CZ_H    = 496    # content zone height(px) — kết thúc trước footer (y=627)

# OPC relationship type cho image embed
_REL_TYPE_IMAGE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
)
# URI cho SVG blip extension (PowerPoint 2016+)
_SVG_EXT_URI = "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"

# Mapping placeholder type integer → tên đọc được
_PH_TYPE_MAP = {
    1:  "title",
    2:  "body",
    3:  "center_title",
    4:  "subtitle",
    5:  "body",        # date
    6:  "footer",
    7:  "slide_number",
    10: "object",
    15: "title",
    None: "body",
}


# ═══════════════════════════════════════════════════════════════════════════
# 1. parse_master_info
# ═══════════════════════════════════════════════════════════════════════════

def parse_master_info(pptx_bytes: bytes) -> dict:
    """
    Đọc file PPTX và trả về thông tin layout, theme, placeholder.

    Args:
        pptx_bytes: Nội dung file .pptx dạng bytes.

    Returns:
        dict với keys: slide_width_emu, slide_height_emu, layouts, theme_colors, default_font
    """
    prs = Presentation(io.BytesIO(pptx_bytes))

    # ── Kích thước slide ────────────────────────────────────────────────────
    slide_width_emu  = int(prs.slide_width)
    slide_height_emu = int(prs.slide_height)

    # ── Theme colors ────────────────────────────────────────────────────────
    theme_colors = _extract_theme_colors(prs)

    # ── Default font ────────────────────────────────────────────────────────
    default_font = _extract_default_font(prs)

    # ── Layouts ─────────────────────────────────────────────────────────────
    layouts = []
    for idx, layout in enumerate(prs.slide_layouts):
        placeholders = []
        for ph in layout.placeholders:
            ph_type_int = ph.placeholder_format.type if ph.placeholder_format else None
            ph_type_str = _PH_TYPE_MAP.get(
                int(ph_type_int) if ph_type_int is not None else None, "body"
            )
            try:
                left   = int(ph.left)   if ph.left   is not None else 0
                top    = int(ph.top)    if ph.top    is not None else 0
                width  = int(ph.width)  if ph.width  is not None else 0
                height = int(ph.height) if ph.height is not None else 0
            except Exception:
                left = top = width = height = 0

            placeholders.append({
                "idx":        ph.placeholder_format.idx if ph.placeholder_format else 0,
                "type":       ph_type_str,
                "left_emu":   left,
                "top_emu":    top,
                "width_emu":  width,
                "height_emu": height,
            })

        layouts.append({
            "index":        idx,
            "name":         layout.name or f"Layout {idx}",
            "placeholders": placeholders,
        })

    return {
        "slide_width_emu":  slide_width_emu,
        "slide_height_emu": slide_height_emu,
        "layouts":          layouts,
        "theme_colors":     theme_colors,
        "default_font":     default_font,
    }


def _extract_theme_colors(prs: Presentation) -> dict:
    """Đọc màu theme từ slide master XML bằng XPath đúng namespace."""
    color_keys = ["dk1", "lt1", "dk2", "lt2",
                  "accent1", "accent2", "accent3", "accent4", "accent5", "accent6"]
    result = {k: "" for k in color_keys}

    try:
        # Lấy theme XML từ slide master đầu tiên
        master = prs.slide_masters[0]
        theme_element = master.theme_color_map  # nếu có
    except Exception:
        return result

    try:
        # Truy cập raw XML của theme qua part
        theme_part = master.part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        )
        tree = etree.fromstring(theme_part.blob)

        for key in color_keys:
            # Tìm node <a:dk1>, <a:lt1>... trong fmtScheme/fontScheme và clrScheme
            nodes = tree.xpath(
                f".//a:clrScheme/a:{key}/*",
                namespaces=_DML_NS
            )
            if nodes:
                node = nodes[0]
                # sysClr có lastClr, srgbClr có val
                color_val = node.get("lastClr") or node.get("val") or ""
                if color_val:
                    result[key] = f"#{color_val.upper()}"
    except Exception:
        pass  # Nếu không đọc được theme → trả về dict rỗng

    return result


def _extract_default_font(prs: Presentation) -> str:
    """Đọc font chữ mặc định từ theme hoặc slide master."""
    try:
        master = prs.slide_masters[0]
        theme_part = master.part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        )
        tree = etree.fromstring(theme_part.blob)
        # Tìm latin font trong majorFont
        nodes = tree.xpath(
            ".//a:fontScheme/a:majorFont/a:latin",
            namespaces=_DML_NS
        )
        if nodes:
            typeface = nodes[0].get("typeface", "")
            if typeface:
                return typeface
        # Fallback: minorFont
        nodes = tree.xpath(
            ".//a:fontScheme/a:minorFont/a:latin",
            namespaces=_DML_NS
        )
        if nodes:
            typeface = nodes[0].get("typeface", "")
            if typeface:
                return typeface
    except Exception:
        pass
    return "Calibri"  # fallback mặc định


# ═══════════════════════════════════════════════════════════════════════════
# 2. find_best_layout
# ═══════════════════════════════════════════════════════════════════════════

# Map data-layout value → danh sách tên layout có thể có (English + Japanese + Korean + Chinese)
# Được dùng bởi find_best_layout() để tìm layout phù hợp theo tên thực tế
_DATA_LAYOUT_TO_NAMES = {
    "title-slide":    [
        "title slide", "title, slide",
        "タイトル スライド", "제목 슬라이드", "标题幻灯片",
    ],
    "section-header": [
        "section header", "section",
        "セクション見出し", "단락머리글", "节标题",
    ],
    "content":        [
        "title and content", "title, content", "content",
        "タイトルとコンテンツ", "标题和内容",
        "제목 및 내용",
    ],
    "two-column":     [
        "two content", "comparison", "two column", "2 content",
        "2 つのコンテンツ", "比較", "两栏内容", "두 내용",
    ],
    "title-only":     [
        "title only", "タイトルのみ",
    ],
    "content-caption": [
        "content with caption", "タイトル付きのコンテンツ",
    ],
    "picture-caption": [
        "picture with caption", "タイトル付きの図",
    ],
    "big-stat":       ["title and content", "content"],
    "blank":          ["blank", "白紙"],
}

# Fallback index khi không tìm thấy bằng tên (clamp đến max_idx khi dùng)
_LAYOUT_FALLBACK_INDEX = {
    "title-slide":    0,
    "section-header": 2,
    "content":        1,
    "two-column":     3,
    "big-stat":       1,
    "blank":          6,
}


def find_best_layout(prs: Presentation, svg_data_layout: str):
    """
    Map data-layout từ SVG sang layout object phù hợp trong master PPTX.

    Tìm kiếm theo tên (có hỗ trợ tiếng Nhật, Hàn, Trung, Anh).
    Fallback theo index cố định.
    Fallback cuối: prs.slide_layouts[0].

    Args:
        prs:             Presentation object (python-pptx).
        svg_data_layout: Giá trị thuộc tính data-layout trong SVG.

    Returns:
        SlideLayout object phù hợp nhất.
    """
    layout_key   = (svg_data_layout or "content").lower().strip()
    target_names = _DATA_LAYOUT_TO_NAMES.get(layout_key,
                   _DATA_LAYOUT_TO_NAMES.get("content", []))
    layouts      = prs.slide_layouts
    max_idx      = len(layouts) - 1

    # 1. Tìm theo tên layout (có hỗ trợ tất cả ngôn ngữ)
    for target in target_names:
        target_low = target.lower()
        for layout in layouts:
            if target_low in (layout.name or "").lower():
                return layout

    # 2. Fallback theo index cố định
    fb_idx = _LAYOUT_FALLBACK_INDEX.get(layout_key, 1)
    fb_idx = min(fb_idx, max_idx) if max_idx >= 0 else 0
    if 0 <= fb_idx <= max_idx:
        return layouts[fb_idx]

    # 3. Fallback cuối cùng
    return layouts[-1] if layouts else layouts[0]


# ═══════════════════════════════════════════════════════════════════════════
# 3. extract_svg_semantic_content
# ═══════════════════════════════════════════════════════════════════════════

def extract_svg_semantic_content(slide_svg: str) -> dict:
    """
    Parse SVG XML để lấy nội dung từ các data-role attribute.
    Hỗ trợ cấu trúc mới: <metadata><slide-layout>, data-type, data-level,
    data-source, data-author. Bỏ qua <g data-role="decorative">.

    Args:
        slide_svg: Chuỗi SVG hoàn chỉnh của một slide đơn.

    Returns:
        dict với layout, title, subtitle, content, content_left, content_right, footer.
    """
    result = {
        "layout":        "content",
        "title":         "",
        "subtitle":      "",
        "content":       [],
        "content_left":  [],
        "content_right": [],
        "footer":        "",
    }

    try:
        parser = etree.XMLParser(recover=True, encoding="utf-8")
        try:
            root = etree.fromstring(slide_svg.encode("utf-8"), parser)
        except Exception:
            return result

        def _localname(el) -> str:
            """
            Lấy tag name không có namespace prefix.
            An toàn với XML Comment/PI nodes (lxml biểu diễn bằng callable tag).
            """
            tag = el.tag
            if callable(tag):          # Comment, ProcessingInstruction, v.v.
                return ""
            if not isinstance(tag, str):
                return ""
            return tag.split("}")[1] if "{" in tag else tag

        def _get_full_text(el) -> str:
            """
            Lấy toàn bộ text content của element (kể cả <tspan> lồng nhau).
            An toàn với XML Comment nodes (callable tag trong lxml).
            """
            parts = []
            for node in el.iter():
                # Bỏ qua Comment/PI nodes (callable tag)
                if callable(node.tag):
                    continue
                if node.get("data-role", "") == "decorative":
                    continue
                if node.text and node.text.strip():
                    parts.append(node.text.strip())
                if node is not el and node.tail and node.tail.strip():
                    parts.append(node.tail.strip())
            return " ".join(parts)

        def _parse_content_items(group_el) -> list:
            """
            Parse các text item trong một content group.
            Chỉ lấy direct children có data-type hoặc là <text>,
            tránh đọc đệ quy trùng lặp qua tspan.
            """
            items = []

            def _process_child(el):
                # Bỏ qua Comment/PI nodes
                if callable(el.tag):
                    return
                lname = _localname(el)
                role  = el.get("data-role", "")

                if role == "decorative":
                    return

                dtype = el.get("data-type", "")

                if dtype:
                    # Element có data-type rõ ràng → lấy toàn bộ text
                    text = _get_full_text(el)
                    if text:
                        try:
                            level = int(el.get("data-level", "1") or 1)
                        except (ValueError, TypeError):
                            level = 1
                        items.append({
                            "text":   text,
                            "type":   dtype,
                            "level":  level,
                            "source": el.get("data-source", ""),
                            "author": el.get("data-author", ""),
                        })
                elif lname == "text":
                    # <text> không có data-type → paragraph
                    text = _get_full_text(el)
                    if text:
                        items.append({
                            "text":   text,
                            "type":   "paragraph",
                            "level":  1,
                            "source": "",
                            "author": "",
                        })
                elif lname == "g" and not role:
                    # <g> không có data-role → đệ quy vào children trực tiếp
                    for child in el:
                        _process_child(child)

            for child in group_el:
                _process_child(child)

            # Fallback: nếu không có items, lấy toàn bộ text của group
            if not items:
                text = _get_full_text(group_el)
                if text:
                    items.append({
                        "text":   text,
                        "type":   "paragraph",
                        "level":  1,
                        "source": "",
                        "author": "",
                    })
            return items

        # ── Tìm slide root ────────────────────────────────────────────────
        slide_g = None
        for el in root.iter():
            # Guard: bỏ qua Comment/PI nodes (tag là callable trong lxml)
            if callable(el.tag):
                continue
            el_id = el.get("id") or ""   # dùng "or" thay default="" vì Comment node
            if el_id and re.match(r"^slide_\d+$", el_id):
                slide_g = el
                break
        if slide_g is None:
            # Fallback: dùng <g> đầu tiên con của root (wrapper <g>)
            for el in root:
                if callable(el.tag):     # bỏ qua Comment nodes
                    continue
                if _localname(el) == "g":
                    slide_g = el
                    break
        if slide_g is None:
            slide_g = root

        # ── Đọc data-layout từ attribute ─────────────────────────────────
        result["layout"] = slide_g.get("data-layout", "content") or "content"

        # ── Đọc <metadata><slide-layout> (override data-layout nếu có) ───
        # Cũng đọc trên root trong trường hợp slide_g = root
        for search_el in [slide_g, root]:
            for child in search_el:
                if callable(child.tag):    # bỏ qua Comment nodes
                    continue
                if _localname(child) == "metadata":
                    for meta_child in child:
                        if _localname(meta_child) == "slide-layout":
                            val = (meta_child.text or "").strip()
                            if val:
                                result["layout"] = val
                    break

        # ── Duyệt direct children của slide_g để tìm data-role groups ────
        for el in slide_g:
            # Bỏ qua Comment/PI nodes (callable tag trong lxml)
            if callable(el.tag):
                continue
            lname = _localname(el)
            role  = el.get("data-role", "").lower().strip()

            if lname == "metadata" or role == "decorative":
                continue

            if role == "title":
                result["title"] = _get_full_text(el)
            elif role == "subtitle":
                result["subtitle"] = _get_full_text(el)
            elif role == "footer":
                result["footer"] = _get_full_text(el)
            elif role == "content":
                result["content"] = _parse_content_items(el)
            elif role == "content-left":
                result["content_left"] = _parse_content_items(el)
            elif role == "content-right":
                result["content_right"] = _parse_content_items(el)

    except Exception:
        pass  # Trả về dict mặc định nếu parse lỗi

    return result


# ═══════════════════════════════════════════════════════════════════════════
# 4. build_pptx_with_master
# ═══════════════════════════════════════════════════════════════════════════

def _find_placeholder(slide, idx: int):
    """Tìm placeholder theo idx. Trả về None nếu không tồn tại."""
    for ph in slide.placeholders:
        try:
            if ph.placeholder_format.idx == idx:
                return ph
        except Exception:
            continue
    return None


def _copy_sldnum_placeholder(slide, layout, slide_number: int) -> None:
    """
    Copy thủ công placeholder số trang từ layout vào slide XML.

    python-pptx KHÔNG tự động copy các footer/date/sldNum placeholder
    từ layout sang slide khi gọi add_slide(). Nếu không có element này
    trực tiếp trong slide XML, PowerPoint hiển thị ô số trang trống.

    Args:
        slide:        Slide object mới tạo.
        layout:       Layout đã dùng cho slide.
        slide_number: Chỉ số 1-based để tạo cNvPr id duy nhất.
    """
    try:
        import copy as _copy
        from lxml import etree as _etree

        P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

        # Tìm <p:sp> có <p:ph type="sldNum"/> trong layout
        sldnum_sp = None
        for sp_ph in layout.placeholders:
            try:
                if sp_ph.placeholder_format.type.name == "SLIDE_NUMBER":
                    sldnum_sp = sp_ph._element
                    break
            except Exception:
                continue

        if sldnum_sp is None:
            return  # Layout không có sldNum placeholder

        # Deep copy — không sửa layout gốc
        sp_copy = _copy.deepcopy(sldnum_sp)

        # Cập nhật cNvPr id → phải là số nguyên duy nhất trong slide
        cNvPr = sp_copy.find(f"{{{P_NS}}}nvSpPr/{{{P_NS}}}cNvPr")
        if cNvPr is not None:
            cNvPr.set("id", str(2000 + slide_number))

        slide.shapes._spTree.append(sp_copy)
    except Exception:
        pass  # Không crash build


# ═══════════════════════════════════════════════════════════════════════════
# 5. SVG crop & embed helpers (cho hybrid master mode)
# ═══════════════════════════════════════════════════════════════════════════

def _crop_svg_to_content_zone(full_slide_svg: str) -> str:
    """
    Tạo SVG chỉ chứa content zone bằng cách đổi viewBox.

    Quy trình:
      1. Parse SVG bằng BeautifulSoup lxml-xml
      2. Xóa <g data-role="title"> (đã được set qua ph[0] placeholder)
      3. Đổi viewBox = "88 131 1104 496" → crop đúng vào content zone
      4. Đổi width/height = 1104/496

    Returns:
        SVG string đã crop. Trả về full_slide_svg nếu lỗi (không crash).
    """
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(full_slide_svg, "lxml-xml")
        svg_tag = soup.find("svg")
        if svg_tag is None:
            return full_slide_svg
        # Xóa title group — sẽ hiển thị qua ph[0] placeholder của master
        for g in svg_tag.find_all("g", attrs={"data-role": "title"}):
            g.decompose()
        # Crop viewBox vào content zone
        svg_tag["viewBox"] = f"{_CZ_X} {_CZ_Y} {_CZ_W} {_CZ_H}"
        svg_tag["width"]   = str(_CZ_W)
        svg_tag["height"]  = str(_CZ_H)
        return str(soup)
    except Exception:
        return full_slide_svg


def _strip_text_roles_from_svg(full_slide_svg: str) -> str:
    """
    Tạo SVG background-only cho title slide.

    Xóa data-role="title" và data-role="subtitle" — các text này
    sẽ hiển thị qua ph[0] và ph[1] placeholder của master.
    Giữ nguyên viewBox full slide (0 0 1280 720) để cover toàn bộ nền.

    Returns:
        SVG string chỉ còn background/decorative.
        Trả về full_slide_svg nếu lỗi (không crash).
    """
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(full_slide_svg, "lxml-xml")
        svg_tag = soup.find("svg")
        if svg_tag is None:
            return full_slide_svg
        for role in ("title", "subtitle"):
            for g in svg_tag.find_all("g", attrs={"data-role": role}):
                g.decompose()
        # Giữ nguyên viewBox full slide
        svg_tag["viewBox"] = f"0 0 {_SVG_W} {_SVG_H}"
        svg_tag["width"]   = str(_SVG_W)
        svg_tag["height"]  = str(_SVG_H)
        return str(soup)
    except Exception:
        return full_slide_svg


def _embed_svg_at(
    slide,
    svg_content: str,
    left_emu: int,
    top_emu: int,
    width_emu: int,
    height_emu: int,
    shape_name: str,
    insert_front: bool = False,
) -> None:
    """
    Embed SVG vào slide tại vị trí và kích thước chỉ định (EMU).

    Args:
        slide:        python-pptx Slide object.
        svg_content:  Chuỗi SVG đã xử lý.
        left_emu:     Tọa độ X từ cạnh trái slide (EMU).
        top_emu:      Tọa độ Y từ cạnh trên slide (EMU).
        width_emu:    Chiều rộng vùng hiển thị (EMU).
        height_emu:   Chiều cao vùng hiển thị (EMU).
        shape_name:   Tên shape trong PPTX (để debug).
        insert_front: True → insert vào đầu spTree (SVG ở layer dưới placeholder).
                      False → append vào cuối (SVG ở layer trên).
    """
    try:
        # ── Correct OPC import paths ─────────────────────────────────────────
        from pptx.opc.package import Part as _OpcPart, PackURI as _PackURI
        from lxml import etree as _etree
        import hashlib, time as _time

        svg_bytes  = svg_content.encode("utf-8")
        slide_part = slide.part

        # Tên part duy nhất để tránh conflict trong OPC package
        uid_str  = f"{shape_name}_{_time.monotonic()}"
        uid_hex  = hashlib.md5(uid_str.encode()).hexdigest()        # 32 hex chars
        # shape id phải là số nguyên dương (OOXML spec)
        shape_id = abs(int(uid_hex[:8], 16)) % 32766 + 2           # 2..32767

        partname = _PackURI(f"/ppt/media/svg_{uid_hex[:12]}.svg")

        # ── Minimal 1×1 white PNG fallback (required by PowerPoint viewers) ──
        # PowerPoint cần blip fallback PNG; nếu không có → hiện tam giác lỗi
        _PNG_1X1 = (
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00"
            b"\x00\x01\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\x0cIDATx"
            b"\x9cc\xf8\xff\xff?\x00\x05\xfe\x02\xfe\xdc\xccY\xe7\x00\x00"
            b"\x00\x00IEND\xaeB`\x82"
        )
        png_partname = _PackURI(f"/ppt/media/png_{uid_hex[:12]}.png")
        png_part = _OpcPart(
            partname     = png_partname,
            content_type = "image/png",
            blob         = _PNG_1X1,
            package      = slide_part.package,
        )
        svg_part = _OpcPart(
            partname     = partname,
            content_type = "image/svg+xml",
            blob         = svg_bytes,
            package      = slide_part.package,
        )
        rId_png = slide_part.relate_to(png_part, _REL_TYPE_IMAGE)
        rId_svg = slide_part.relate_to(svg_part, _REL_TYPE_IMAGE)

        # Namespace shortcuts
        PPTX_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
        DML_NS  = "http://schemas.openxmlformats.org/drawingml/2006/main"
        REL_NS  = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
        SVG_NS  = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"

        _p  = lambda t: f"{{{PPTX_NS}}}{t}"
        _a  = lambda t: f"{{{DML_NS}}}{t}"
        _r  = lambda t: f"{{{REL_NS}}}{t}"
        _sv = lambda t: f"{{{SVG_NS}}}{t}"

        pic = _etree.Element(_p("pic"),
                             nsmap={"p": PPTX_NS, "a": DML_NS,
                                    "r": REL_NS, "asvg": SVG_NS})

        # nvPicPr — id phải là số nguyên hợp lệ
        nvPicPr  = _etree.SubElement(pic, _p("nvPicPr"))
        cNvPr    = _etree.SubElement(nvPicPr, _p("cNvPr"))
        cNvPr.set("id", str(shape_id))
        cNvPr.set("name", shape_name)
        cNvPicPr = _etree.SubElement(nvPicPr, _p("cNvPicPr"))
        cNvPicPr.set("preferRelativeResize", "0")
        picLocks = _etree.SubElement(cNvPicPr, _a("picLocks"))
        picLocks.set("noChangeAspect", "1")
        _etree.SubElement(nvPicPr, _p("nvPr"))

        # blipFill: PNG fallback + SVG extension
        blipFill = _etree.SubElement(pic, _p("blipFill"))
        blip     = _etree.SubElement(blipFill, _a("blip"))
        blip.set(_r("embed"), rId_png)          # fallback PNG
        extLst   = _etree.SubElement(blip, _a("extLst"))
        ext      = _etree.SubElement(extLst, _a("ext"))
        ext.set("uri", _SVG_EXT_URI)
        svgBlip  = _etree.SubElement(ext, _sv("svgBlip"))
        svgBlip.set(_r("embed"), rId_svg)       # SVG vector
        stretch  = _etree.SubElement(blipFill, _a("stretch"))
        _etree.SubElement(stretch, _a("fillRect"))

        # spPr — vị trí và kích thước
        spPr     = _etree.SubElement(pic, _p("spPr"))
        xfrm     = _etree.SubElement(spPr, _a("xfrm"))
        off      = _etree.SubElement(xfrm, _a("off"))
        off.set("x", str(int(left_emu)))
        off.set("y", str(int(top_emu)))
        ext_sz   = _etree.SubElement(xfrm, _a("ext"))
        ext_sz.set("cx", str(int(width_emu)))
        ext_sz.set("cy", str(int(height_emu)))
        prstGeom = _etree.SubElement(spPr, _a("prstGeom"))
        prstGeom.set("prst", "rect")
        _etree.SubElement(prstGeom, _a("avLst"))

        # Thêm vào spTree
        sp_tree = slide.shapes._spTree
        if insert_front:
            sp_tree.insert(2, pic)   # index 2: sau nvGrpSpPr và grpSpPr
        else:
            sp_tree.append(pic)

    except Exception as _e:
        import sys
        print(f"[_embed_svg_at] ERROR {shape_name}: {_e}", file=sys.stderr)


def build_pptx_with_master(
    slides: list[dict],
    pptx_bytes: bytes,
    master_info: dict,
) -> io.BytesIO:
    """
    Tạo PPTX mới dựa trên master từ pptx_bytes bằng cách map nội dung SVG
    vào đúng placeholder của từng layout.

    Args:
        slides:      Danh sách slide dict từ extract_slides_from_svg().
                     Mỗi item gồm: 'id', 'index', 'svg'.
        pptx_bytes:  Nội dung file master .pptx dạng bytes.
        master_info: Kết quả từ parse_master_info().

    Returns:
        io.BytesIO chứa file PPTX đã map xong.
    """
    if not slides:
        raise ValueError("Danh sách slides không được rỗng.")

    # ── Mở PPTX master ──────────────────────────────────────────────────────
    prs = Presentation(io.BytesIO(pptx_bytes))
    # KHÔNG thay đổi slide_width/slide_height — kế thừa từ master

    # ── Xóa toàn bộ slide hiện có (proper cleanup) ──────────────────────────
    # Phải xóa cả OPC relationship lẫn sldId element, nếu không PPTX bị corrupt
    # vì còn dangling rel trỏ tới slide part đã bị orphan.
    _REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    sldIdLst = prs.slides._sldIdLst
    for sldId_el in list(sldIdLst):
        rId = sldId_el.get(f"{{{_REL_NS}}}id")
        if rId:
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass
        try:
            sldIdLst.remove(sldId_el)
        except Exception:
            pass

    # ── Tính scale EMU/px từ master thực tế (không hardcode) ─────────────────
    _w_emu = int(prs.slide_width)
    _h_emu = int(prs.slide_height)
    _scale = _w_emu / _SVG_W   # EMU per SVG pixel, uniform cả 2 chiều

    # Content zone trong EMU (tính động từ scale)
    _cz_left   = round(_CZ_X * _scale)
    _cz_top    = round(_CZ_Y * _scale)
    _cz_width  = round(_CZ_W * _scale)
    _cz_height = round(_CZ_H * _scale)

    # ── Thêm slide mới theo nội dung SVG ─────────────────────────────────────
    for slide_idx, slide_data in enumerate(slides):
        svg_text   = slide_data.get("svg", "")

        # 1. Parse semantic để lấy title, subtitle, layout
        semantic   = extract_svg_semantic_content(svg_text)
        raw_layout = (
            slide_data.get("data_layout", "")
            or semantic.get("layout", "")
            or "content"
        )
        is_title_slide = raw_layout in ("title-slide", "title slide")

        # 2. Tìm layout phù hợp trong master
        layout = find_best_layout(prs, raw_layout)

        # 3. Thêm slide mới
        slide = prs.slides.add_slide(layout)

        # Copy sldNum placeholder từ layout vào slide → số trang hiển thị đúng
        _copy_sldnum_placeholder(slide, layout, slide_idx + 1)

        # 4. Set title qua ph[0] placeholder → dùng font/style của master
        title_text = semantic.get("title", "")
        _ph0 = _find_placeholder(slide, 0)
        if _ph0 is not None:
            try:
                _ph0.text = title_text
            except Exception:
                pass

        # 5. Xử lý body theo loại layout ──────────────────────────────────────
        if is_title_slide:
            # TITLE SLIDE:
            #   ph[0] đã có title (bước 4)
            #   ph[1] subtitle → set text qua placeholder (master font/style)
            #   Background → embed full SVG đã xóa title/subtitle text
            subtitle_text = semantic.get("subtitle", "")
            _ph1 = _find_placeholder(slide, 1)
            if _ph1 is not None:
                try:
                    _ph1.text = subtitle_text
                except Exception:
                    pass

            # Embed SVG background: full slide, chỉ giữ decorative shapes
            bg_svg = _strip_text_roles_from_svg(svg_text)
            if bg_svg.strip():
                _embed_svg_at(
                    slide        = slide,
                    svg_content  = bg_svg,
                    left_emu     = 0,
                    top_emu      = 0,
                    width_emu    = _w_emu,
                    height_emu   = _h_emu,
                    shape_name   = f"SVGBg_{slide_idx + 1}",
                    insert_front = True,   # dưới placeholder để chữ hiện lên trên
                )

        else:
            # CONTENT SLIDE:
            #   ph[0] đã có title (bước 4)
            #   Body → embed SVG crop vào content zone (giữ NGUYÊN visual design)
            #   KHÔNG tạo textbox, KHÔNG phân rã text
            content_svg = _crop_svg_to_content_zone(svg_text)
            if content_svg.strip():
                _embed_svg_at(
                    slide        = slide,
                    svg_content  = content_svg,
                    left_emu     = _cz_left,
                    top_emu      = _cz_top,
                    width_emu    = _cz_width,
                    height_emu   = _cz_height,
                    shape_name   = f"SVGContent_{slide_idx + 1}",
                    insert_front = False,  # sau placeholder, hiển thị trên cùng
                )
        # ph[12] slide number: PowerPoint tự điền — không cần xử lý

    # ── Lưu ra BytesIO ──────────────────────────────────────────────────────
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output
