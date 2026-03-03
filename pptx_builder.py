"""
pptx_builder.py
---------------
Module tạo file PowerPoint (.pptx) từ danh sách các chuỗi SVG.

Điểm nổi bật:
  - Nhúng SVG trực tiếp vào PPTX (không qua PNG).
  - PowerPoint 2019+ có thể dùng "Convert to Shape" để chỉnh sửa từng element.
  - Sử dụng API nội bộ của python-pptx để thêm part SVG vào package OPC.
"""

import io
from lxml import etree

from pptx import Presentation
from pptx.util import Emu
# python-pptx 1.0+: Part và PackURI đều nằm trong pptx.opc.package
from pptx.opc.package import Part, PackURI


# ─── Kích thước slide 16:9 tính bằng EMU ────────────────────────────────────
# 1 inch = 914400 EMU  →  10 inch × 5.625 inch = 16:9
SLIDE_WIDTH_EMU  = Emu(9144000)   # 10 inch
SLIDE_HEIGHT_EMU = Emu(5143500)   # 5.625 inch

# ─── Namespace dùng trong XML của PowerPoint ────────────────────────────────
NSMAP = {
    "a":    "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p":    "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r":    "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "asvg": "http://schemas.microsoft.com/office/drawing/2016/SVG/main",
}

# Relationship type cho ảnh trong slide
REL_TYPE_IMAGE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
)

# Content-type cho file SVG
SVG_CONTENT_TYPE = "image/svg+xml"

# URI {}  cho extension SVG blip trong DrawingML
SVG_EXTENSION_URI = "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"


def _build_pic_xml(rId: str, shape_id: int, slide_name: str) -> etree._Element:
    """
    Xây dựng phần tử XML <p:pic> để hiển thị ảnh SVG trên slide.

    Args:
        rId       : Relationship ID trỏ tới file SVG trong package.
        shape_id  : ID duy nhất của shape trong slide.
        slide_name: Tên hiển thị của shape.

    Returns:
        Phần tử lxml <p:pic> hoàn chỉnh.
    """
    # Khai báo namespace đầy đủ cho phần tử gốc
    PPTX_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    DML_NS  = "http://schemas.openxmlformats.org/drawingml/2006/main"
    REL_NS  = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    SVG_NS  = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"

    # Hàm tiện ích tạo tag có namespace
    def p(tag):  return f"{{{PPTX_NS}}}{tag}"
    def a(tag):  return f"{{{DML_NS}}}{tag}"
    def r(attr): return f"{{{REL_NS}}}{attr}"
    def asvg(tag): return f"{{{SVG_NS}}}{tag}"

    # ── <p:pic> ──────────────────────────────────────────────────────────────
    pic = etree.Element(p("pic"), nsmap={
        "p": PPTX_NS, "a": DML_NS, "r": REL_NS, "asvg": SVG_NS
    })

    # ── <p:nvPicPr> — thuộc tính phi hình ảnh ─────────────────────────────
    nvPicPr = etree.SubElement(pic, p("nvPicPr"))

    cNvPr = etree.SubElement(nvPicPr, p("cNvPr"))
    cNvPr.set("id",   str(shape_id))
    cNvPr.set("name", slide_name)

    cNvPicPr = etree.SubElement(nvPicPr, p("cNvPicPr"))
    cNvPicPr.set("preferRelativeResize", "0")
    picLocks = etree.SubElement(cNvPicPr, a("picLocks"))
    picLocks.set("noChangeAspect", "1")

    etree.SubElement(nvPicPr, p("nvPr"))  # phần tử rỗng bắt buộc

    # ── <p:blipFill> — nguồn ảnh ──────────────────────────────────────────
    blipFill = etree.SubElement(pic, p("blipFill"))

    blip = etree.SubElement(blipFill, a("blip"))
    blip.set(r("embed"), rId)

    # Extension dành cho SVG (PowerPoint 2016/2019+)
    extLst = etree.SubElement(blip, a("extLst"))
    ext = etree.SubElement(extLst, a("ext"))
    ext.set("uri", SVG_EXTENSION_URI)
    svgBlip = etree.SubElement(ext, asvg("svgBlip"))
    svgBlip.set(r("embed"), rId)

    stretch = etree.SubElement(blipFill, a("stretch"))
    etree.SubElement(stretch, a("fillRect"))  # vừa khít slide

    # ── <p:spPr> — vị trí & kích thước ────────────────────────────────────
    spPr = etree.SubElement(pic, p("spPr"))

    xfrm = etree.SubElement(spPr, a("xfrm"))
    off = etree.SubElement(xfrm, a("off"))
    off.set("x", "0")
    off.set("y", "0")
    ext_size = etree.SubElement(xfrm, a("ext"))
    ext_size.set("cx", str(int(SLIDE_WIDTH_EMU)))
    ext_size.set("cy", str(int(SLIDE_HEIGHT_EMU)))

    prstGeom = etree.SubElement(spPr, a("prstGeom"))
    prstGeom.set("prst", "rect")
    etree.SubElement(prstGeom, a("avLst"))

    return pic


def _add_svg_to_slide(slide, svg_content: str, slide_index: int) -> None:
    """
    Nhúng một chuỗi SVG vào slide PowerPoint.

    Quy trình:
      1. Mã hóa SVG thành bytes.
      2. Tạo OPC Part mới với content-type 'image/svg+xml'.
      3. Tạo relationship từ slide → part SVG.
      4. Thêm phần tử <p:pic> vào spTree của slide.

    Args:
        slide       : Đối tượng slide của python-pptx.
        svg_content : Chuỗi SVG cho slide này.
        slide_index : Số thứ tự slide (0-based), dùng để đặt tên file.
    """
    svg_bytes = svg_content.encode("utf-8")
    slide_part = slide.part

    # Đường dẫn nội bộ của file SVG trong package (ví dụ: /ppt/media/svg_slide_1.svg)
    svg_partname = PackURI(f"/ppt/media/svg_slide_{slide_index + 1}.svg")

    # Tạo Part OPC chứa dữ liệu SVG
    # Lưu ý: python-pptx 1.0+ đổi thứ tự tham số: Part(partname, content_type, package, blob)
    svg_part = Part(
        partname=svg_partname,
        content_type=SVG_CONTENT_TYPE,
        package=slide_part.package,
        blob=svg_bytes,
    )

    # Thêm relationship: slide → SVG part → trả về rId
    rId = slide_part.relate_to(svg_part, REL_TYPE_IMAGE)

    # Xây dựng XML <p:pic> và gắn vào spTree của slide
    shape_id   = slide_index + 2            # id 1 thường là background
    shape_name = f"SVG Slide {slide_index + 1}"
    pic_element = _build_pic_xml(rId, shape_id, shape_name)

    sp_tree = slide.shapes._spTree               # lxml element chứa tất cả shapes
    sp_tree.append(pic_element)


def build_pptx_from_slides(slides: list[dict]) -> io.BytesIO:
    """
    Tạo file PPTX từ danh sách các slide SVG.

    Args:
        slides: Danh sách dict từ svg_processor.extract_slides_from_svg().
                Mỗi dict gồm 'id', 'index', 'svg'.

    Returns:
        io.BytesIO chứa dữ liệu file PPTX, sẵn sàng để gửi về client.

    Raises:
        ValueError: Nếu danh sách slides rỗng.
    """
    if not slides:
        raise ValueError("Danh sách slides không được rỗng.")

    # ── Tạo Presentation mới với tỉ lệ 16:9 ──────────────────────────────
    prs = Presentation()
    prs.slide_width  = SLIDE_WIDTH_EMU
    prs.slide_height = SLIDE_HEIGHT_EMU

    # Layout trống (blank) — index 6 trong hầu hết các template mặc định
    blank_layout = prs.slide_layouts[6]

    for slide_data in slides:
        # Thêm slide mới với layout trống
        slide = prs.slides.add_slide(blank_layout)

        # Nhúng SVG vào slide
        _add_svg_to_slide(
            slide=slide,
            svg_content=slide_data["svg"],
            slide_index=slide_data["index"] - 1,   # chuyển sang 0-based
        )

    # ── Lưu vào bộ nhớ và trả về ──────────────────────────────────────────
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def build_pptx_from_slides_with_master(
    slides: list[dict],
    pptx_bytes: bytes,
    master_info: dict,
) -> io.BytesIO:
    """
    Wrapper gọi master_handler.build_pptx_with_master().
    Tách biệt để app.py không cần import trực tiếp master_handler.

    Args:
        slides:      Danh sách slide dict từ extract_slides_from_svg().
        pptx_bytes:  Nội dung file master .pptx dạng bytes.
        master_info: Kết quả từ parse_master_info().

    Returns:
        io.BytesIO chứa file PPTX đã map vào master.
    """
    from master_handler import build_pptx_with_master
    return build_pptx_with_master(slides, pptx_bytes, master_info)
